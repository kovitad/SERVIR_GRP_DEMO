from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE

OUT = Path(__file__).with_name("GRP_AI_Observability_and_Evaluation_Proposal.pptx")
_logo_candidates = [
    Path(__file__).parents[1] / "assets" / "servir-global-collaborative.png",
    Path(__file__).parents[2] / "public" / "assets" / "servir-global-collaborative.png",
]
LOGO = next((path for path in _logo_candidates if path.exists()), _logo_candidates[0])

# SERVIR presentation palette
BLUE = RGBColor(35, 128, 176)
LIGHT_BLUE = RGBColor(74, 154, 218)
DARK = RGBColor(33, 33, 33)
MID = RGBColor(89, 89, 89)
GREY = RGBColor(109, 110, 113)
PALE = RGBColor(241, 247, 250)
LINE = RGBColor(215, 225, 230)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(145, 175, 61)
AMBER = RGBColor(255, 171, 64)
RED = RGBColor(201, 79, 60)
PURPLE = RGBColor(112, 64, 144)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def fill(shape, color, transparency=0):
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.fill.transparency = transparency
    if hasattr(shape, "line"): shape.line.fill.background()


def line(shape, color=LINE, width=1.0, dash=None):
    shape.line.color.rgb = color; shape.line.width = Pt(width)
    if dash is not None: shape.line.dash_style = dash


def text(slide, x, y, w, h, value, size=18, color=DARK, bold=False,
         font="Roboto", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
         margin=0, linesp=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin); tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = linesp
    r = p.add_run(); r.text = value; r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    return box


def rich_lines(slide, x, y, w, h, items, size=15, bullet_color=BLUE, gap=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(.03); tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        if isinstance(item, tuple): heading, body = item
        else: heading, body = "", item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.level = 0
        rb = p.add_run(); rb.text = "●  "; rb.font.size = Pt(size-2); rb.font.color.rgb = bullet_color
        if heading:
            rh = p.add_run(); rh.text = heading; rh.font.name = "Roboto"; rh.font.size = Pt(size); rh.font.bold = True; rh.font.color.rgb = DARK
        rr = p.add_run(); rr.text = body; rr.font.name = "Roboto"; rr.font.size = Pt(size); rr.font.color.rgb = MID
    return box


def rect(slide, x, y, w, h, color=WHITE, radius=True, border=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    fill(sh, color)
    if border: line(sh, border, 1)
    return sh


def pill(slide, x, y, w, label, color=BLUE, text_color=WHITE):
    sh = rect(slide, x, y, w, .34, color)
    text(slide, x, y+.01, w, .30, label, 9, text_color, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def header(slide, title, kicker=None, num=None):
    if kicker: text(slide, .62, .30, 8.5, .25, kicker.upper(), 9, BLUE, True, font="Roboto Condensed")
    text(slide, .62, .57, 10.7, .52, title, 25, DARK, True, font="Roboto Condensed")
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.62), Inches(1.15), Inches(1.0), Inches(.045)).fill.solid()
    bar = slide.shapes[-1]; bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    if LOGO.exists(): slide.shapes.add_picture(str(LOGO), Inches(10.92), Inches(.30), width=Inches(1.78))
    if num is not None: text(slide, 12.42, 7.08, .32, .22, str(num), 8, GREY, align=PP_ALIGN.RIGHT)
    text(slide, .62, 7.08, 5.2, .22, "SERVIR GLOBAL RISK PLATFORM · PROPOSAL", 8, GREY, True, font="Roboto Condensed")


def arrow(slide, x1, y1, x2, y2, color=BLUE, width=2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line(c, color, width); c.line.end_arrowhead = True
    return c


def node(slide, x, y, w, h, label, sub="", color=WHITE, accent=BLUE):
    rect(slide, x, y, w, h, color, border=LINE)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(.07), Inches(h)).fill.solid()
    a=slide.shapes[-1]; a.fill.fore_color.rgb=accent; a.line.fill.background()
    text(slide, x+.18, y+.13, w-.28, .28, label, 13, DARK, True, font="Roboto Condensed")
    if sub: text(slide, x+.18, y+.45, w-.28, h-.53, sub, 9.5, MID)


# 1 — Title
s=prs.slides.add_slide(blank)
fill(s.background, DARK)
s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(.12)).fill.solid(); s.shapes[-1].fill.fore_color.rgb=BLUE; s.shapes[-1].line.fill.background()
rect(s,.62,.48,2.25,.58,WHITE)
if LOGO.exists(): s.shapes.add_picture(str(LOGO), Inches(.78), Inches(.64), width=Inches(1.9))
pill(s,10.65,.55,1.98,"TEAM PROPOSAL",BLUE)
text(s,.72,1.65,10.8,.45,"AI OBSERVABILITY",16,LIGHT_BLUE,True,font="Roboto Condensed")
text(s,.72,2.08,11.6,1.20,"Making every GRP answer\ntraceable, testable and governable",35,WHITE,True,font="Roboto Condensed")
text(s,.75,3.63,9.8,.54,"LangGraph workflow tracing · Langfuse Cloud · repeatable evaluation",18,RGBColor(210,225,233))
rect(s,.72,4.70,11.85,1.25,RGBColor(47,47,47),border=RGBColor(75,75,75))
text(s,1.00,4.96,11.25,.66,"PROOF POINT",9,AMBER,True,font="Roboto Condensed")
text(s,1.00,5.25,11.1,.42,"Can we explain why an answer was produced—and know whether the next release is better?",19,WHITE,True,font="Roboto Condensed")
text(s,.72,6.88,6,.25,"SERVIR GLOBAL RISK PLATFORM",9,RGBColor(180,190,194),True,font="Roboto Condensed")

# 2 — Problem
s=prs.slides.add_slide(blank); header(s,"A completed request can still be wrong","Problem",2)
text(s,.72,1.45,5.3,.62,"A GRP answer crosses multiple components",20,DARK,True,font="Roboto Condensed")
components=[("Agent workflow","routing and decisions",BLUE),("LLM","reasoning and generation",PURPLE),("MCP / APIs","tools and external services",AMBER),("Datasets","source and version",GREEN),("GIS calculations","method and outputs",LIGHT_BLUE)]
for i,(a,b,c) in enumerate(components):
    y=2.12+i*.80; node(s,.72,y,3.25,.62,a,b,WHITE,c)
    if i<4: arrow(s,4.04,y+.31,4.48,y+.31,GREY,1.4)
# right blind spot panel
rect(s,4.62,1.50,8.02,4.88,PALE,border=LINE)
text(s,4.98,1.83,7.3,.36,"Normal application logs usually answer:",15,MID,False)
pill(s,9.68,1.78,2.38,"REQUEST COMPLETED",GREEN)
text(s,4.98,2.35,7.0,.38,"But they do not reliably explain:",18,DARK,True,font="Roboto Condensed")
rich_lines(s,4.98,2.85,7.05,3.15,[
    ("Why ","the agent produced that answer"),
    ("Which ","data, model, prompt and tools it used"),
    ("Where ","an incorrect result originated"),
    ("Whether ","a new version performs better or worse"),
    ("How much ","each request costs and how long it takes")
],15,RED,10)

# 3 — Solution scope
s=prs.slides.add_slide(blank); header(s,"Observe the complete path from question to evidence","Solution",3)
text(s,.72,1.43,11.6,.38,"Request-level observability connects workflow behaviour, evidence provenance and quality evaluation.",17,MID)
cols=[
("TRACE","End-to-end workflow","LangGraph steps and decisions\nMCP/API inputs and outputs\nErrors and latency",BLUE),
("VERSION","Prompts, models and data","Prompt and model versions\nDataset and calculation provenance\nTool configuration",PURPLE),
("MEASURE","Usage and quality","Tokens, latency and cost\nGolden-question scores\nUser feedback and alerts",GREEN),
("GOVERN","Access and audit","Masking and retention\nRole-based access\nInvestigation history",AMBER)]
for i,(k,t,b,c) in enumerate(cols):
    x=.72+i*3.03; rect(s,x,2.08,2.78,3.52,WHITE,border=LINE)
    pill(s,x+.22,2.32,1.00,k,c)
    text(s,x+.22,2.83,2.32,.63,t,17,DARK,True,font="Roboto Condensed")
    for j,linev in enumerate(b.split("\n")):
        text(s,x+.22,3.65+j*.55,2.35,.42,"• "+linev,11,MID)
rect(s,.72,5.92,11.86,.62,PALE,border=LINE)
text(s,.92,6.06,11.45,.30,"One trace ID links the user request, workflow steps, generations, tools, evidence and evaluation scores.",14,BLUE,True,align=PP_ALIGN.CENTER)

# 4 — Trace anatomy
s=prs.slides.add_slide(blank); header(s,"What one observable GRP request looks like","Trace anatomy",4)
text(s,.72,1.38,10.8,.42,"Example: “Why prioritise populations farther than 1 km from candidate shelters?”",17,DARK,True,font="Roboto Condensed")
steps=[("1","VALIDATE","AOI · scenario · access",.6,BLUE),("2","RETRIEVE","evidence IDs · versions",2.95,GREEN),("3","CALCULATE","method · parameters",5.30,LIGHT_BLUE),("4","GENERATE","prompt · model · tokens",7.65,PURPLE),("5","CHECK","citations · limitations",10.00,AMBER)]
for i,(n,t,b,x,c) in enumerate(steps):
    rect(s,x,2.28,2.02,1.25,WHITE,border=LINE)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+.13),Inches(2.45),Inches(.40),Inches(.40)); fill(circ,c)
    text(s,x+.13,2.48,.40,.28,n,10,WHITE,True,align=PP_ALIGN.CENTER)
    text(s,x+.64,2.39,1.20,.25,t,11,c,True,font="Roboto Condensed")
    text(s,x+.14,2.87,1.72,.36,b,9.5,MID,align=PP_ALIGN.CENTER)
    if i<4: arrow(s,x+2.05,2.90,x+2.30,2.90,GREY,1.5)
rect(s,.72,4.05,11.86,1.78,PALE,border=LINE)
text(s,.98,4.30,2.4,.28,"LANGFUSE TRACE",11,BLUE,True,font="Roboto Condensed")
metrics=[("2.8 s","total latency"),("1,420","input tokens"),("286","output tokens"),("$0.006","estimated cost"),("3","evidence records"),("PASS","citation check")]
for i,(v,lbl) in enumerate(metrics):
    x=.98+i*1.88
    text(s,x,4.75,1.55,.38,v,20,DARK,True,font="Roboto Condensed",align=PP_ALIGN.CENTER)
    text(s,x,5.17,1.55,.28,lbl,9,GREY,align=PP_ALIGN.CENTER)
text(s,.78,6.14,11.8,.36,"Observation-first: evaluate the logical root for the final answer, retrieval observations for relevance, and generation observations for groundedness.",10,BLUE,True,align=PP_ALIGN.CENTER)

# 5 — Evaluator design
s=prs.slides.add_slide(blank); header(s,"Use the cheapest reliable evaluator for each failure mode","Evaluation design",5)
text(s,.72,1.40,11.4,.40,"Do not ask one “God Evaluator” to produce a vague quality score. Split checks so every failure tells us what to fix.",16,MID)
# deterministic panel
rect(s,.72,2.02,5.68,3.72,WHITE,border=LINE)
pill(s,1.00,2.29,1.55,"CODE CHECKS",GREEN)
text(s,1.00,2.82,4.95,.46,"Deterministic first",20,DARK,True,font="Roboto Condensed")
text(s,1.00,3.25,4.95,.34,"Exact · fast · cheap · repeatable",11,GREEN,True)
rich_lines(s,1.00,3.77,4.92,1.60,[
("Figures match: ","assessment values equal the expected output"),
("Evidence exists: ","every citation ID and version resolves"),
("Schema valid: ","required fields and types are present"),
("Limitations present: ","mandatory safety statements are included")
],11,GREEN,4)
# judge panel
rect(s,6.67,2.02,5.91,3.72,PALE,border=LINE)
pill(s,6.95,2.29,1.68,"LLM JUDGES",PURPLE)
text(s,6.95,2.82,5.02,.46,"Narrow judgment only",20,DARK,True,font="Roboto Condensed")
text(s,6.95,3.25,5.02,.34,"One criterion · reasoning first · verdict last",11,PURPLE,True)
rich_lines(s,6.95,3.77,5.05,1.50,[
("Groundedness: ","are material claims supported by context?"),
("Question relevance: ","does the answer address the planner’s need?"),
("Action usefulness: ","are proposed actions relevant and usable?")
],11,PURPLE,5)
rect(s,.72,6.02,11.86,.52,PALE,border=LINE)
text(s,.94,6.13,11.42,.28,"Preferred verdict:  pass  |  fail  |  unknown     ·     Grade outcomes in evidence—not claims in prose.",12,BLUE,True,align=PP_ALIGN.CENTER)

# 6 — Calibration and release loop
s=prs.slides.add_slide(blank); header(s,"Calibrate the judge before trusting the score","Evaluator calibration",6)
cal=[("1","LABEL","10–20 real cases"),("2","DEFINE","one precise criterion"),("3","PROMPT","context + ignore rules"),("4","TEST","held-out labels"),("5","COMPARE","agreement by class"),("6","ITERATE","review disagreements")]
for i,(n,a,b) in enumerate(cal):
    x=.42+i*2.13
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+.62),Inches(1.72),Inches(.58),Inches(.58)); fill(circ,[BLUE,GREEN,PURPLE,LIGHT_BLUE,AMBER,BLUE][i])
    text(s,x+.62,1.84,.58,.24,n,10,WHITE,True,align=PP_ALIGN.CENTER)
    text(s,x,2.48,1.82,.28,a,11,DARK,True,font="Roboto Condensed",align=PP_ALIGN.CENTER)
    text(s,x,2.84,1.82,.32,b,9,MID,align=PP_ALIGN.CENTER)
    if i<5: arrow(s,x+1.72,2.02,x+2.04,2.02,LINE,1.5)
rect(s,.72,3.55,5.66,2.12,WHITE,border=LINE)
pill(s,1.00,3.82,1.55,"OFFLINE",BLUE)
text(s,1.00,4.30,4.95,.38,"Experiments for release decisions",17,DARK,True,font="Roboto Condensed")
text(s,1.00,4.79,4.95,.56,"Run the same golden questions against baseline and candidate prompt, model or workflow versions.",12,MID)
rect(s,6.67,3.55,5.91,2.12,PALE,border=LINE)
pill(s,6.95,3.82,1.55,"ONLINE",GREEN)
text(s,6.95,4.30,5.05,.38,"Observation-level monitoring",17,DARK,True,font="Roboto Condensed")
text(s,6.95,4.79,5.05,.56,"Evaluate targeted root, retrieval or generation observations; sample live traffic to control cost.",12,MID)
rect(s,.72,5.95,11.86,.58,RGBColor(255,245,231),border=AMBER)
text(s,.94,6.09,11.40,.28,"Warning: 90% overall agreement can hide a judge that always says “pass”. Validate pass, fail and unknown separately.",11,RED,True,align=PP_ALIGN.CENTER)

# 7 — Use case
s=prs.slides.add_slide(blank); header(s,"One simple use case in the current prototype","Demonstration",7)
rect(s,.72,1.48,11.86,.72,PALE,border=LINE)
text(s,1.00,1.66,11.30,.36,'“Why should areas farther than 1 km from candidate shelters be prioritised?”',18,DARK,True,font="Roboto Condensed",align=PP_ALIGN.CENTER)
flow=[("APP RESULT","18,640 exposed\n11,240 beyond 1 km",BLUE),("RETRIEVE","3 governed\nevidence records",GREEN),("OPENAI","planning explanation\nwith citations",PURPLE),("EVALUATE","code checks + narrow\njudge verdicts",AMBER)]
for i,(a,b,c) in enumerate(flow):
    x=.78+i*3.02; node(s,x,2.70,2.58,1.12,a,b,WHITE,c)
    if i<3: arrow(s,x+2.61,3.25,x+2.92,3.25,GREY,1.5)
text(s,.72,4.30,4.9,.36,"Answer contract",18,DARK,True,font="Roboto Condensed")
rich_lines(s,.72,4.72,5.35,1.80,["Short planning explanation","Reported figures with evidence IDs","Recommended follow-up actions","Explicit mocked-data and shelter limitations"],12,BLUE,5)
rect(s,6.55,4.25,6.03,1.98,WHITE,border=LINE)
pill(s,6.82,4.52,1.98,"EVALUATION STACK",GREEN)
text(s,6.82,5.01,5.40,.34,"4 deterministic checks + 3 narrow judges",15,DARK,True,font="Roboto Condensed")
text(s,6.82,5.43,5.40,.54,"Each score attaches to the relevant observation and explains exactly which failure mode passed, failed or was unknown.",11,MID)

# 8 — Architecture
s=prs.slides.add_slide(blank); header(s,"Minimal architecture for the proof of concept","Implementation",8)
arch=[("BROWSER","Existing GRP prototype","Question + assessment context",BLUE),
      ("BACKEND API","Node.js service","Keys remain server-side",DARK),
      ("LANGGRAPH","Five-node workflow","Validate · retrieve · generate · check",PURPLE),
      ("OPENAI","LLM provider","Answer + token usage",GREEN)]
for i,(a,b,c,d) in enumerate(arch):
    x=.55+i*3.08; node(s,x,1.72,2.62,1.20,a,b+"\n"+c,WHITE,d)
    if i<3: arrow(s,x+2.65,2.32,x+2.98,2.32,GREY,1.5)
# Langfuse below
rect(s,3.62,3.52,6.08,1.30,PALE,border=BLUE)
text(s,3.90,3.77,5.5,.30,"LANGFUSE CLOUD",13,BLUE,True,font="Roboto Condensed",align=PP_ALIGN.CENTER)
text(s,3.90,4.15,5.5,.40,"Root + operation observations · tokens · cost · evaluator scores",12,DARK,True,align=PP_ALIGN.CENTER)
arrow(s,6.66,3.44,6.66,2.98,BLUE,2)
rect(s,.72,5.33,11.86,.85,WHITE,border=LINE)
text(s,.98,5.54,2.15,.28,"PRIVATE ENVIRONMENT",10,RED,True,font="Roboto Condensed")
text(s,3.10,5.46,8.90,.42,"OPENAI_API_KEY  ·  LANGFUSE_PUBLIC_KEY  ·  LANGFUSE_SECRET_KEY  ·  LANGFUSE_BASE_URL",12,MID,True)
text(s,.85,6.44,11.55,.26,"No API keys in browser JavaScript or GitHub. No LangGraph API key is required.",11,GREY,align=PP_ALIGN.CENTER)

# 9 — Live demo
s=prs.slides.add_slide(blank); header(s,"A five-minute team demonstration","Demo flow",9)
steps=[
("01","Run assessment","Select Phaya Thai · RP100"),
("02","Ask AI","Generate planning explanation"),
("03","Inspect answer","Figures · citations · limitations"),
("04","Open Langfuse","Root · retrieval · generation observations"),
("05","Compare runs","Code checks · judge verdicts · cost")]
for i,(n,a,b) in enumerate(steps):
    y=1.45+i*.94
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(.78),Inches(y),Inches(.54),Inches(.54)); fill(circ,BLUE if i<4 else GREEN)
    text(s,.78,y+.08,.54,.27,n,10,WHITE,True,align=PP_ALIGN.CENTER)
    text(s,1.55,y-.01,3.35,.29,a,15,DARK,True,font="Roboto Condensed")
    text(s,1.55,y+.34,3.75,.30,b,10,MID)
    if i<4: arrow(s,1.05,y+.62,1.05,y+.87,LINE,2)
rect(s,6.00,1.48,6.58,4.56,PALE,border=LINE)
text(s,6.35,1.82,5.85,.34,"What the team should see",18,DARK,True,font="Roboto Condensed")
rich_lines(s,6.35,2.35,5.55,3.20,[
("Traceability: ","every answer links to its workflow and evidence"),
("Diagnosis: ","the failing component is visible"),
("Evaluation: ","quality changes are measured, not assumed"),
("Economics: ","latency, tokens and cost are explicit"),
("Governance: ","access, masking and retention can be controlled")
],14,BLUE,10)

# 10 — Proof and measures
s=prs.slides.add_slide(blank); header(s,"What we want to prove","Success criteria",10)
proofs=[
("REPRODUCE","A reviewer can identify the exact prompt, model, evidence and tool versions behind an answer.",BLUE,"100% traced requests"),
("DIAGNOSE","A failed golden question can be traced to retrieval, calculation, tool or generation.",RED,"Cause found < 10 min"),
("COMPARE","Two workflow versions can be compared with the same calibrated question set and evaluators.",PURPLE,"No silent regression"),
("CONTROL","The team can see and set expectations for latency, tokens and cost.",GREEN,"Cost/request visible"),
("GOVERN","Sensitive content is masked and traces follow access and retention policies.",AMBER,"Policy checks pass")]
for i,(a,b,c,d) in enumerate(proofs):
    y=1.42+i*1.02; rect(s,.72,y,11.86,.78,WHITE,border=LINE)
    pill(s,.94,y+.21,1.35,a,c)
    text(s,2.50,y+.17,7.65,.42,b,12,DARK,True,font="Roboto Condensed")
    text(s,10.34,y+.23,1.88,.28,d,10,c,True,align=PP_ALIGN.CENTER)

# 11 — benefits and next step
s=prs.slides.add_slide(blank); header(s,"Start small, prove value, then expand","Recommendation",11)
text(s,.72,1.42,5.45,.35,"Expected benefits",19,DARK,True,font="Roboto Condensed")
rich_lines(s,.72,1.87,5.45,3.85,[
"Diagnose incorrect answers faster",
"Reproduce results and investigate failures",
"Prevent quality regression between releases",
"Demonstrate evidence behind reported figures",
"Control latency, token usage and operating cost",
"Compare tools, prompts and models objectively",
"Improve stakeholder confidence and governance"
],13,BLUE,7)
rect(s,6.45,1.42,6.13,4.45,PALE,border=LINE)
pill(s,6.78,1.73,1.72,"PROPOSED PILOT",GREEN)
text(s,6.78,2.24,5.38,.62,"One observable planning explanation",21,DARK,True,font="Roboto Condensed")
rich_lines(s,6.78,3.00,5.15,1.85,[
("Scope: ","one question, one workflow, 5–10 evidence records"),
("Evaluation: ","4 code checks, 3 judges, 10–20 labeled cases"),
("Stack: ","LangGraph + OpenAI + Langfuse Cloud"),
("Effort: ","2–4 development days plus a team labeling session"),
("Decision: ","expand only after the trace and evaluation value is demonstrated")
],12,GREEN,7)
rect(s,6.78,5.20,5.42,.43,BLUE)
text(s,6.78,5.27,5.42,.25,"NEXT: CREATE LANGFUSE PROJECT + PRIVATE KEYS",10,WHITE,True,font="Roboto Condensed",align=PP_ALIGN.CENTER)
text(s,.72,6.42,11.86,.36,"The goal is not more logging. The goal is evidence that the AI system behaves as intended.",17,BLUE,True,font="Roboto Condensed",align=PP_ALIGN.CENTER)

# 12 — Design basis
s=prs.slides.add_slide(blank); header(s,"Design basis and reference material","References",12)
text(s,.72,1.45,11.35,.42,"The proposed evaluation pattern follows Langfuse’s observation-first and calibrated-evaluator guidance.",17,DARK,True,font="Roboto Condensed")
refs=[
("LLM-as-a-Judge","Observation-level evaluators for live traffic; experiments for controlled comparisons; numeric, categorical and boolean scores.","https://langfuse.com/docs/evaluation/evaluation-methods/llm-as-a-judge"),
("Writing good evaluators","Prefer deterministic checks; one evaluator per failure mode; binary/categorical verdicts; label real cases and calibrate judges.","https://langfuse.com/academy/evaluate/writing-evaluators"),
("Implementation note","Trace-level evaluators are deprecated for Langfuse v4. The pilot should target logical root and operation-level observations.","Langfuse documentation reviewed for this proposal")]
for i,(a,b,u) in enumerate(refs):
    y=2.05+i*1.30; rect(s,.72,y,11.86,1.02,WHITE,border=LINE)
    pill(s,.98,y+.22,1.82,"REFERENCE "+str(i+1),[BLUE,PURPLE,AMBER][i])
    text(s,3.02,y+.16,3.05,.29,a,15,DARK,True,font="Roboto Condensed")
    text(s,3.02,y+.48,8.92,.28,b,10.5,MID)
    text(s,3.02,y+.76,8.92,.20,u,8,BLUE)
rect(s,.72,6.20,11.86,.45,PALE,border=LINE)
text(s,.92,6.30,11.46,.24,"Human review remains the calibration reference. An LLM judge is another AI component and must itself be measured.",11,RED,True,align=PP_ALIGN.CENTER)

# Core metadata
prs.core_properties.title = "GRP AI Observability and Evaluation Proposal"
prs.core_properties.subject = "LangGraph, OpenAI and Langfuse Cloud proof of concept"
prs.core_properties.author = "SERVIR Global Risk Platform"
prs.core_properties.keywords = "AI observability, evaluation, LangGraph, Langfuse, OpenAI, GRP"
prs.save(OUT)
print(OUT)
